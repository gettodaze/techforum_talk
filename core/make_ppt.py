import typing as tp

from pathlib import Path

from pptx import Presentation

from core.config import Config


INPUT_FILE = Config.INPUT / "presentation.md"
OUTPUT_FILE = Config.OUTPUT / "presentation.pptx"


class PowerPoint:
    def __init__(self) -> None:
        self.prs = Presentation()

    def add_lines(self, lines: tp.List[str]) -> None:
        code = False
        bullets: tp.List[str] = []
        title = ""
        for line in lines:
            line = line.strip("\n")
            flag = line[: line.find(" ")] if " " in line else line

            if flag == "CODE":
                code = True
            elif flag == "END":
                self.add_bullet_slide(bullets=bullets, title=title)
                bullets = []
                title = ""
                code = False
            elif not code and flag == "#":
                self.add_bullet_slide(title=title, bullets=bullets)
                bullets = []
                title = line[2:]
            elif not code and flag == "##":
                self.add_bullet_slide(title=title, bullets=bullets)
                bullets = []
                title = line[3:]
            else:
                bullets.append(line)

    def save(self, fp: str) -> str:
        self.prs.save(fp)
        return fp

    def add_bullet_slide(self, title: str = "", bullets: tp.Iterable[str] = ()):
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        shapes.title.text = title

        tf = shapes.placeholders[1].text_frame
        # tf.text = 'Find the bullet slide layout'
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0


if __name__ == "__main__":
    lines = INPUT_FILE.read_text().splitlines()
    ppt = PowerPoint()
    ppt.add_lines(lines)
    ppt.save(OUTPUT_FILE)
    print(OUTPUT_FILE)
